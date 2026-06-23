#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
G  = RGBColor(0x2E,0xA8,0x4F)  # green
N  = RGBColor(0x1B,0x2B,0x6B)  # navy
W  = RGBColor(0xFF,0xFF,0xFF)  # white
LB = RGBColor(0xEB,0xF5,0xFC)  # light blue bg
GR = RGBColor(0xCC,0xCC,0xCC)  # gray
DK = RGBColor(0x33,0x33,0x33)  # dark text
R  = RGBColor(0xC0,0x39,0x2B)  # red
OR = RGBColor(0xE6,0x7E,0x22)  # orange
BL = RGBColor(0x29,0x80,0xB9)  # blue
PR = RGBColor(0x8E,0x44,0xAD)  # purple
GG = RGBColor(0x27,0xAE,0x60)  # green2

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLK = prs.slide_layouts[6]

NAV = ["Contexte Général","Objectifs du PFA",
       "Architecture & Technologies","Analyse & Conception",
       "Conclusion & Perspectives"]

FH = Inches(0.48)
NH = Inches(0.40)
ML = Inches(0.65)

# ── primitives ──────────────────────────────────────
def rct(sl,x,y,w,h,fc=None,lc=None,lw=None):
    s=sl.shapes.add_shape(1,x,y,w,h)
    if fc: s.fill.solid(); s.fill.fore_color.rgb=fc
    else: s.fill.background()
    if lc: s.line.color.rgb=lc; s.line.width=lw or Pt(0.75)
    else: s.line.fill.background()
    return s

def ovl(sl,x,y,w,h,fc):
    s=sl.shapes.add_shape(9,x,y,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fc
    s.line.fill.background()
    return s

def txt(sl,x,y,w,h,text,sz=12,bold=False,col=None,al=PP_ALIGN.LEFT,wrap=True):
    if col is None: col=DK
    b=sl.shapes.add_textbox(x,y,w,h)
    tf=b.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=al
    rn=p.add_run(); rn.text=text
    rn.font.size=Pt(sz); rn.font.bold=bold; rn.font.color.rgb=col
    return b

def shape_txt(sl,x,y,w,h,text,sz=12,bold=False,col=None,al=PP_ALIGN.CENTER,fc=G,lc=None):
    """Rectangle shape with centered text"""
    if col is None: col=W
    s=sl.shapes.add_shape(1,x,y,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fc
    if lc: s.line.color.rgb=lc; s.line.width=Pt(0.75)
    else: s.line.fill.background()
    tf=s.text_frame; tf.word_wrap=True
    try: tf._txBody.get_or_add_bodyPr().set('anchor','ctr')
    except: pass
    p=tf.paragraphs[0]; p.alignment=al
    rn=p.add_run(); rn.text=text
    rn.font.size=Pt(sz); rn.font.bold=bold; rn.font.color.rgb=col
    return s

def footer(sl,pg):
    rct(sl,0,SH-FH,SW,FH,fc=G)
    txt(sl,ML,SH-FH+Inches(0.07),Inches(5),FH,"SecureWatch",13,True,W)
    txt(sl,SW-Inches(1.3),SH-FH+Inches(0.07),Inches(1.1),FH,str(pg),17,True,W,PP_ALIGN.RIGHT)

def nav_bar(sl,active):
    iw=SW/5
    for i,lbl in enumerate(NAV):
        fc=G if i==active else W
        tc=W if i==active else N
        s=sl.shapes.add_shape(1,i*iw,0,iw,NH)
        s.fill.solid(); s.fill.fore_color.rgb=fc
        s.line.color.rgb=GR; s.line.width=Pt(0.5)
        tf=s.text_frame; tf.word_wrap=False
        try: tf._txBody.get_or_add_bodyPr().set('anchor','ctr')
        except: pass
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        rn=p.add_run(); rn.text=lbl; rn.font.size=Pt(9); rn.font.bold=True; rn.font.color.rgb=tc

def sec_title(sl,title):
    y0=NH+Inches(0.15)
    txt(sl,ML,y0,SW-2*ML,Inches(0.42),title,20,True,N)
    rct(sl,ML,y0+Inches(0.44),Inches(0.6),Inches(0.055),fc=G)
    return y0+Inches(0.57)

def bullet(sl,x,y,w,btxt,rtxt,sz=13):
    ovl(sl,x,y+Inches(0.12),Inches(0.09),Inches(0.09),G)
    b=sl.shapes.add_textbox(x+Inches(0.18),y,w-Inches(0.18),Inches(0.56))
    tf=b.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if btxt:
        r1=p.add_run(); r1.text=btxt+" "; r1.font.size=Pt(sz); r1.font.bold=True; r1.font.color.rgb=G
    r2=p.add_run(); r2.text=rtxt; r2.font.size=Pt(sz); r2.font.color.rgb=DK

def cards_3x2(sl,items,y0):
    cw=(SW-2*ML-Inches(0.4))/3
    ch=Inches(1.75)
    gx=Inches(0.2); gy=Inches(0.15)
    for i,(icon,title,desc) in enumerate(items):
        c=i%3; row=i//3
        cx=ML+c*(cw+gx); cy=y0+row*(ch+gy)
        rct(sl,cx,cy,cw,ch,fc=LB,lc=GR)
        txt(sl,cx,cy+Inches(0.06),cw,Inches(0.65),icon,28,False,N,PP_ALIGN.CENTER)
        txt(sl,cx+Inches(0.08),cy+Inches(0.72),cw-Inches(0.16),Inches(0.32),title,10,True,N,PP_ALIGN.CENTER)
        txt(sl,cx+Inches(0.08),cy+Inches(1.04),cw-Inches(0.16),ch-Inches(1.08),desc,9,False,DK,PP_ALIGN.CENTER,True)

def cards_4(sl,items,y0):
    cw=(SW-2*ML-Inches(0.6))/4
    gap=Inches(0.2)
    for i,(icon,title,desc) in enumerate(items):
        cx=ML+i*(cw+gap)
        txt(sl,cx,y0,cw,Inches(1.05),icon,52,False,N,PP_ALIGN.CENTER)
        txt(sl,cx,y0+Inches(1.1),cw,Inches(0.38),title,10,True,G,PP_ALIGN.CENTER)
        txt(sl,cx,y0+Inches(1.52),cw,Inches(1.1),desc,10,False,DK,PP_ALIGN.CENTER,True)

def sec_intro(num,l1,l2,icon,pg):
    sl=prs.slides.add_slide(BLK)
    rct(sl,0,0,SW,SH-FH,fc=LB)
    txt(sl,ML,Inches(1.5),Inches(2.2),Inches(3.5),str(num),100,True,N,PP_ALIGN.LEFT)
    txt(sl,ML+Inches(2.0),Inches(2.0),Inches(5.5),Inches(1.2),l1,46,True,N)
    if l2: txt(sl,ML+Inches(2.0),Inches(3.2),Inches(5.5),Inches(1.1),l2,46,True,N)
    uly=Inches(4.35) if l2 else Inches(3.3)
    rct(sl,ML+Inches(2.0),uly,Inches(1.1),Inches(0.07),fc=G)
    rct(sl,SW*0.58,Inches(1.6),Inches(0.03),Inches(3.2),fc=GR)
    txt(sl,SW*0.62,Inches(1.2),Inches(4),Inches(4),icon,88,False,N,PP_ALIGN.CENTER)
    footer(sl,pg)

# ══════════════════════════════════════════════════════
# S1 — PAGE DE GARDE
# ══════════════════════════════════════════════════════
s1=prs.slides.add_slide(BLK)
rct(s1,0,SH*0.55,SW,SH*0.45,fc=LB)
rct(s1,Inches(4.2),Inches(0.35),Inches(0.9),Inches(0.9),fc=N)
txt(s1,Inches(4.2),Inches(0.35),Inches(0.9),Inches(0.9),"EMSI",12,True,W,PP_ALIGN.CENTER)
txt(s1,Inches(5.2),Inches(0.38),Inches(5.0),Inches(0.42),"ECOLE MAROCAINE DES SCIENCES DE L'INGÉNIEUR",11,True,N)
txt(s1,Inches(5.2),Inches(0.78),Inches(5.5),Inches(0.25),"Membre de  •  HONORIS UNITED UNIVERSITIES",9,False,GR)
txt(s1,0,Inches(1.55),SW,Inches(0.55),"PROJET DE FIN D'ANNÉE",24,True,N,PP_ALIGN.CENTER)
txt(s1,0,Inches(2.08),SW,Inches(0.4),"2025 / 2026",16,False,RGBColor(0x77,0x77,0x77),PP_ALIGN.CENTER)
rct(s1,Inches(2.0),Inches(2.7),Inches(9.333),Inches(1.05),fc=G)
txt(s1,Inches(2.0),Inches(2.7),Inches(9.333),Inches(1.05),
    "SecureWatch : Centralisation des Logs &\nDétection d'Incidents de Sécurité",
    21,True,W,PP_ALIGN.CENTER)
rct(s1,Inches(2.0),Inches(3.95),Inches(9.333),Inches(0.02),fc=GR)
txt(s1,Inches(2.5),Inches(4.2),Inches(3.5),Inches(0.38),"RÉALISÉ PAR",13,True,N)
txt(s1,Inches(2.5),Inches(4.62),Inches(3.5),Inches(0.38),"AGRAD Hicham",13,False,DK)
rct(s1,SW*0.5,Inches(4.1),Inches(0.02),Inches(1.6),fc=GR)
txt(s1,Inches(7.0),Inches(4.2),Inches(5.5),Inches(0.38),"MEMBRES DU JURY",13,True,N)
txt(s1,Inches(7.0),Inches(4.62),Inches(5.8),Inches(0.9),
    "Pr. ________________________   (Encadrant)\nPr. ________________________",13,False,DK)
rct(s1,0,SH-Inches(0.12),SW,Inches(0.12),fc=G)

# ══════════════════════════════════════════════════════
# S2 — PLAN
# ══════════════════════════════════════════════════════
s2=prs.slides.add_slide(BLK)
txt(s2,ML,Inches(0.5),Inches(3),Inches(0.7),"PLAN",40,True,N)
rct(s2,ML,Inches(1.22),Inches(0.6),Inches(0.07),fc=G)
icons_p=["🔍","🎯","🖥","🧩","✅"]
labels_p=["Contexte\nGénéral","Objectifs\ndu PFA","Architecture &\nTechnologies","Analyse &\nConception","Conclusion &\nPerspectives"]
aw=(SW-2*ML-Inches(0.4))/5; ah=Inches(1.2); ay=Inches(2.6)
for i in range(5):
    ax=ML+i*(aw+Inches(0.1))
    txt(s2,ax,ay-Inches(0.48),aw,Inches(0.42),str(i+1),22,True,N,PP_ALIGN.CENTER)
    try:
        shp=s2.shapes.add_shape(13,ax,ay,aw,ah)
        shp.fill.solid(); shp.fill.fore_color.rgb=G; shp.line.fill.background()
    except:
        rct(s2,ax,ay,aw,ah,fc=G)
    txt(s2,ax,ay,aw,ah,icons_p[i],34,False,W,PP_ALIGN.CENTER)
    txt(s2,ax,ay+ah+Inches(0.1),aw,Inches(0.7),labels_p[i],11,False,N,PP_ALIGN.CENTER,True)
footer(s2,2)

# ══════════════════════════════════════════════════════
# S3 — Section 1 intro
# ══════════════════════════════════════════════════════
sec_intro(1,"Contexte","Général","🔍",3)

# ══════════════════════════════════════════════════════
# S4 — Problématique
# ══════════════════════════════════════════════════════
s4=prs.slides.add_slide(BLK)
nav_bar(s4,0)
y4=sec_title(s4,"Problématique")
bi4=[
    ("Multiplication des cyberattaques :","brute force, injections SQL, scanners offensifs, anomalies géographiques — menaces de plus en plus diversifiées."),
    ("Logs dispersés :","dans une architecture distribuée, les traces d'attaques sont éparpillées entre plusieurs services, difficiles à corréler."),
    ("Détection réactive :","sans monitoring centralisé, les incidents sont détectés trop tard, augmentant leur impact sur les systèmes."),
    ("Besoin d'approche proactive :","centraliser, analyser et alerter en temps réel est devenu indispensable pour toute infrastructure moderne."),
]
by4=y4+Inches(0.05)
for bt,rt in bi4:
    bullet(s4,ML,by4,Inches(6.8),bt,rt,13); by4+=Inches(0.67)
tags=[("🔑 Brute Force",R),("💉 Injection SQL",PR),("🔍 Scanners",OR),
      ("🌍 Anomalie Géo",GG),("🚫 Rate Limit",N),("⛔ Accès 403/401",BL)]
tx=SW-Inches(4.9); ty=y4+Inches(0.05)
txt(s4,tx,ty,Inches(4.5),Inches(0.38),"⚠  Types d'attaques ciblées",13,True,N)
ty+=Inches(0.48)
for i,(lbl,col) in enumerate(tags):
    tag_x=tx+(i%2)*Inches(2.2)
    rct(s4,tag_x,ty,Inches(2.0),Inches(0.38),fc=col)
    txt(s4,tag_x,ty,Inches(2.0),Inches(0.38),lbl,9,True,W,PP_ALIGN.CENTER)
    if i%2==1: ty+=Inches(0.52)
footer(s4,4)

# ══════════════════════════════════════════════════════
# S5 — Présentation SecureWatch
# ══════════════════════════════════════════════════════
s5=prs.slides.add_slide(BLK)
nav_bar(s5,0)
y5=sec_title(s5,"Présentation du projet SecureWatch")
lw5=Inches(5.8)
rct(s5,ML,y5+Inches(0.1),Inches(0.75),Inches(0.75),fc=N)
txt(s5,ML,y5+Inches(0.1),Inches(0.75),Inches(0.75),"🛡",24,False,W,PP_ALIGN.CENTER)
txt(s5,ML+Inches(0.85),y5+Inches(0.12),Inches(3.5),Inches(0.38),"SecureWatch",22,True,N)
txt(s5,ML+Inches(0.85),y5+Inches(0.52),Inches(4),Inches(0.3),"Plateforme de Surveillance de Sécurité",11,False,GR)
desc5=("SecureWatch est une plateforme de surveillance de sécurité en temps réel, développée dans le cadre "
       "du PFA 2025-2026. Elle centralise la collecte des logs, détecte automatiquement les incidents de sécurité "
       "et génère des alertes en temps réel dans une architecture microservices conteneurisée (14 conteneurs Docker).")
b5=s5.shapes.add_textbox(ML,y5+Inches(0.95),lw5,Inches(1.35))
tf5=b5.text_frame; tf5.word_wrap=True
rn5=tf5.paragraphs[0].add_run(); rn5.text=desc5; rn5.font.size=Pt(12); rn5.font.color.rgb=DK
bdgs=["14 conteneurs Docker","ELK Stack 8.12.0","Prometheus + Grafana","9 scénarios d'attaque","RBAC 3 rôles"]
bx5=ML; by5=y5+Inches(2.42)
for bdg in bdgs:
    bw5=Inches(0.13)*len(bdg)+Inches(0.3)
    rct(s5,bx5,by5,bw5,Inches(0.32),fc=LB,lc=GR)
    txt(s5,bx5+Inches(0.08),by5,bw5,Inches(0.32),bdg,9,True,N)
    bx5+=bw5+Inches(0.12)
    if bx5>Inches(6.3): bx5=ML; by5+=Inches(0.38)
rx5=ML+lw5+Inches(0.4); rw5=SW-rx5-ML; rh5=SH-FH-y5-Inches(0.15)
rct(s5,rx5,y5+Inches(0.1),rw5,rh5,fc=LB,lc=GR)
txt(s5,rx5,y5+Inches(0.1),rw5,rh5,"📸\n[dashboard_overview.png]",12,False,GR,PP_ALIGN.CENTER)
footer(s5,5)

# ══════════════════════════════════════════════════════
# S6 — Menaces ciblées
# ══════════════════════════════════════════════════════
s6=prs.slides.add_slide(BLK)
nav_bar(s6,0)
y6=sec_title(s6,"Menaces & types d'attaques ciblés")
cards_3x2(s6,[
    ("🔑","Brute Force","≥5 échecs de connexion depuis la même IP\n→ détection CRITIQUE + counter Prometheus"),
    ("💉","Injection SQL","Regex sur champs login → blocage 400\n+ log SQL_INJECTION (CRITICAL)"),
    ("🔍","Scanners Offensifs","Détection via User-Agent : sqlmap, Nikto,\nNessus, masscan, nuclei, gobuster"),
    ("🌍","Anomalie Géographique","Requêtes depuis IPs à risque\n(RU, CN, KP, IR, TOR, NG) → GEO_ANOMALY"),
    ("🚫","Accès Non Autorisés","Réponses 401/403 agrégées par Logstash\n→ UNAUTHORIZED / FORBIDDEN"),
    ("⚡","Dépassement de Débit","Rate limiting Nginx 3 zones\n(5r/m auth, 30r/m API) → 429 RATE_LIMITED"),
],y6+Inches(0.1))
footer(s6,6)

# ══════════════════════════════════════════════════════
# S7 — Section 2 intro
# ══════════════════════════════════════════════════════
sec_intro(2,"Objectifs","du PFA","🎯",7)

# ══════════════════════════════════════════════════════
# S8 — Objectifs du PFA
# ══════════════════════════════════════════════════════
s8=prs.slides.add_slide(BLK)
nav_bar(s8,1)
y8=sec_title(s8,"Objectifs du PFA")
cards_3x2(s8,[
    ("📚","Centraliser les logs","ELK Stack (Elasticsearch + Logstash\n+ Kibana + Filebeat)"),
    ("🎯","Détecter les attaques","en temps réel — 9 types d'incidents,\n10 event_types Logstash"),
    ("🔔","Alerter automatiquement","Prometheus + Alertmanager\n11 règles d'alerte configurées"),
    ("📊","Visualiser les données","Dashboard SPA + Grafana + Kibana\ncarte géographique D3.js"),
    ("🔒","Sécuriser les accès","JWT HS256, TLS/HTTPS, rate limiting\nRBAC 3 rôles (Admin/Operator/User)"),
    ("🐛","Simuler les attaques","9 scénarios automatisés\n(brute force, SQLi, scanner, géo…)"),
],y8+Inches(0.08))
footer(s8,8)

# ══════════════════════════════════════════════════════
# S9 — Section 3 intro
# ══════════════════════════════════════════════════════
sec_intro(3,"Architecture &","Technologies","🖥",9)

# ══════════════════════════════════════════════════════
# S10 — Technologies utilisées
# ══════════════════════════════════════════════════════
s10=prs.slides.add_slide(BLK)
nav_bar(s10,2)
y10=sec_title(s10,"Technologies utilisées")
techs=[
    ("🐍","Python / Flask",RGBColor(0x37,0x76,0xAB)),
    ("🔎","ELK Stack 8.12",RGBColor(0x00,0xBF,0xB3)),
    ("📈","Prometheus + Grafana",RGBColor(0xE6,0x52,0x2C)),
    ("🐳","Docker Compose",RGBColor(0x24,0x96,0xED)),
    ("🛡","Nginx API Gateway",RGBColor(0x00,0x96,0x39)),
    ("🗄","Redis",RGBColor(0xDC,0x38,0x2D)),
]
tw10=(SW-2*ML-Inches(0.8))/3; th10=Inches(1.6); tgap=Inches(0.4)
for i,(icon,name,col) in enumerate(techs):
    ci=i%3; ri=i//3
    tx10=ML+ci*(tw10+tgap); ty10=y10+Inches(0.15)+ri*(th10+Inches(0.2))
    cz=Inches(0.9); cx10=tx10+tw10/2-cz/2
    s10c=s10.shapes.add_shape(9,cx10,ty10,cz,cz)
    s10c.fill.solid(); s10c.fill.fore_color.rgb=col; s10c.line.fill.background()
    txt(s10,tx10,ty10,tw10,cz,icon,28,False,W,PP_ALIGN.CENTER)
    txt(s10,tx10,ty10+cz+Inches(0.1),tw10,Inches(0.38),name,12,True,N,PP_ALIGN.CENTER)
footer(s10,10)

# ══════════════════════════════════════════════════════
# S11 — Architecture globale
# ══════════════════════════════════════════════════════
s11=prs.slides.add_slide(BLK)
nav_bar(s11,2)
y11=sec_title(s11,"Architecture globale — 14 conteneurs Docker")
lw11=Inches(6.8); lh=Inches(0.57); ly=y11+Inches(0.1); gap11=Inches(0.1)
layers=[
    ("Frontend",    [("SPA Dashboard :3000",G),("Chart.js",N),("D3.js GeoMap",N)]),
    ("API Gateway", [("Nginx :8443 HTTPS/TLS",G),("Rate Limiting 3 zones",N),("JSON Logs",N)]),
    ("Microservices",[("auth-service :5001",G),("api-service :5002",G),("Redis",N)]),
    ("ELK Stack",   [("Filebeat",G),("Logstash :5044",G),("Elasticsearch :9200",N),("Kibana :5601",N)]),
    ("Monitoring",  [("Prometheus :9090",G),("Alertmanager :9093",N),("Grafana :3001",G),("Node Exporter",N)]),
]
for lname,litems in layers:
    lblw=Inches(1.5)
    rct(s11,ML,ly,lblw,lh,fc=N)
    txt(s11,ML,ly,lblw,lh,lname,9,True,W,PP_ALIGN.CENTER)
    ix=ML+lblw+Inches(0.08); iaw=lw11-lblw-Inches(0.08)
    iw=iaw/len(litems)
    for j,(iname,icol) in enumerate(litems):
        ifc=RGBColor(0xE8,0xF8,0xEE) if icol==G else LB
        ilc=G if icol==G else GR
        rct(s11,ix+j*iw,ly+Inches(0.06),iw-Inches(0.04),lh-Inches(0.12),fc=ifc,lc=ilc,lw=Pt(0.75))
        txt(s11,ix+j*iw,ly+Inches(0.06),iw-Inches(0.04),lh-Inches(0.12),iname,8,True,icol,PP_ALIGN.CENTER)
    if lname!="Monitoring":
        txt(s11,ML+lw11/2-Inches(0.3),ly+lh,Inches(0.6),gap11,"↓",10,False,GR,PP_ALIGN.CENTER)
    ly+=lh+gap11
rx11=ML+lw11+Inches(0.4); rw11=SW-rx11-ML; rh11=SH-FH-y11-Inches(0.15)
rct(s11,rx11,y11+Inches(0.1),rw11,rh11,fc=LB,lc=GR)
txt(s11,rx11,y11+Inches(0.1),rw11,rh11,"📸\n[architecture_diagram.png]",11,False,GR,PP_ALIGN.CENTER)
footer(s11,11)

# ══════════════════════════════════════════════════════
# S12 — Section 4 intro
# ══════════════════════════════════════════════════════
sec_intro(4,"Analyse &","Conception","🧩",12)

# ══════════════════════════════════════════════════════
# S13 — Identification des acteurs
# ══════════════════════════════════════════════════════
s13=prs.slides.add_slide(BLK)
nav_bar(s13,3)
y13=sec_title(s13,"Identification des acteurs")
actors=[
    ("👮","Administrateur",N,["Accès complet au dashboard","Gestion des utilisateurs","Vue carte géographique","Accès Kibana / Grafana"]),
    ("⚙️","Opérateur",N,["Consultation des logs de sécurité","Vue carte géographique","Gestion des alertes"]),
    ("👤","Utilisateur",N,["Accès basique au dashboard","Consultation des statistiques"]),
    ("🥷","Attaquant (externe)",R,["Simulé par generate_attacks.py","9 scénarios d'attaque"]),
]
aw13=(SW-2*ML-Inches(0.6))/4; ax13=ML; ay13=y13+Inches(0.15)
for icon,name,col,perms in actors:
    txt(s13,ax13,ay13,aw13,Inches(1.0),icon,50,False,col,PP_ALIGN.CENTER)
    txt(s13,ax13,ay13+Inches(1.02),aw13,Inches(0.4),name,12,True,col,PP_ALIGN.CENTER)
    py=ay13+Inches(1.5)
    for perm in perms:
        ovl(s13,ax13+Inches(0.12),py+Inches(0.1),Inches(0.08),Inches(0.08),G)
        txt(s13,ax13+Inches(0.26),py,aw13-Inches(0.3),Inches(0.38),perm,10,False,DK)
        py+=Inches(0.4)
    ax13+=aw13+Inches(0.2)
footer(s13,13)

# ══════════════════════════════════════════════════════
# S14 — Diagramme de cas d'utilisation
# ══════════════════════════════════════════════════════
s14=prs.slides.add_slide(BLK)
nav_bar(s14,3)
y14=sec_title(s14,"Diagramme de cas d'utilisation")
rct(s14,Inches(2.4),y14+Inches(0.1),Inches(8.2),Inches(4.3),lc=N,fc=LB)
txt(s14,Inches(2.4),y14+Inches(0.1),Inches(8.2),Inches(0.48),"SecureWatch — Système de surveillance",12,True,N,PP_ALIGN.CENTER)
la=[("👮","Admin"),("⚙️","Opérateur"),("👤","User")]
lax14=ML; lay14=y14+Inches(0.5)
for ico,nm in la:
    txt(s14,lax14,lay14,Inches(1.8),Inches(0.55),ico,30,False,N,PP_ALIGN.CENTER)
    txt(s14,lax14,lay14+Inches(0.55),Inches(1.8),Inches(0.3),nm,11,True,N,PP_ALIGN.CENTER)
    lay14+=Inches(1.35)
ucs=["S'authentifier (JWT)","Consulter le dashboard","Voir les logs de sécurité","Gérer les alertes","Voir la carte géographique","Consulter Kibana","Consulter Grafana","Gérer les utilisateurs"]
uc_y=y14+Inches(0.65)
for i,uc in enumerate(ucs):
    uc_x=Inches(2.7) if i%2==0 else Inches(6.9)
    uy=uc_y+(i//2)*Inches(0.82)
    rct(s14,uc_x,uy,Inches(3.9),Inches(0.45),lc=BL)
    txt(s14,uc_x,uy,Inches(3.9),Inches(0.45),uc,10,False,BL,PP_ALIGN.CENTER)
txt(s14,SW-Inches(2.4),y14+Inches(0.4),Inches(1.9),Inches(0.55),"🥷",30,False,R,PP_ALIGN.CENTER)
txt(s14,SW-Inches(2.4),y14+Inches(0.95),Inches(1.9),Inches(0.3),"Attaquant",11,True,R,PP_ALIGN.CENTER)
att_y=y14+Inches(1.35)
for ac in ["Lancer brute force","Injection SQL","Scanner les ports"]:
    rct(s14,SW-Inches(2.4),att_y,Inches(1.9),Inches(0.4),lc=R)
    txt(s14,SW-Inches(2.4),att_y,Inches(1.9),Inches(0.4),ac,9,False,R,PP_ALIGN.CENTER)
    att_y+=Inches(0.52)
footer(s14,14)

# ══════════════════════════════════════════════════════
# S15 — Besoins fonctionnels
# ══════════════════════════════════════════════════════
s15=prs.slides.add_slide(BLK)
nav_bar(s15,3)
y15=sec_title(s15,"Besoins fonctionnels")
bf=[
    ("Centralisation des logs :","Filebeat → Logstash → Elasticsearch (index security-logs-YYYY.MM.dd), 10 event_types configurés"),
    ("Détection en temps réel :","9 types d'incidents — brute force, SQLi, scanners UA, anomalies géo, 403/401, 429, upload suspect, erreurs 500"),
    ("Alertes automatiques :","11 règles Prometheus + Alertmanager, webhook → api-service, SSE push vers dashboard (Redis persistant)"),
    ("Dashboard de monitoring :","Statistiques en temps réel, graphiques Chart.js, carte géographique D3.js + TopoJSON, 7 stat cards"),
    ("Authentification & RBAC :","JWT HS256 (1h), 3 rôles (Admin/Operator/User) avec contrôle d'accès granulaire sur toutes les routes"),
    ("Simulation d'attaques :","9 scénarios automatisés par generate_attacks.py pour tests et démonstrations (brute force, SQLi, scanner…)"),
]
by15=y15+Inches(0.05)
for bt,rt in bf:
    bullet(s15,ML,by15,SW-2*ML,bt,rt,12.5); by15+=Inches(0.62)
footer(s15,15)

# ══════════════════════════════════════════════════════
# S16 — Besoins non-fonctionnels
# ══════════════════════════════════════════════════════
s16=prs.slides.add_slide(BLK)
nav_bar(s16,3)
y16=sec_title(s16,"Besoins non-fonctionnels")
bnf=[
    ("Sécurité :","TLS/HTTPS auto-signé, JWT HS256, rate limiting 3 zones Nginx, sanitize_input() contre log injection, X-Real-IP anti-spoofing"),
    ("Performance :","Traitement des logs en temps réel, alertes SSE push sans polling, réponse dashboard < 1 seconde"),
    ("Disponibilité :","14 conteneurs Docker avec health checks, Redis pour persistance des alertes (survie aux redémarrages)"),
    ("Scalabilité :","Architecture microservices indépendante — chaque service peut être mis à l'échelle séparément"),
    ("Maintenabilité :","Pipeline Logstash configurable, règles Prometheus externalisées en YAML, rapport LaTeX + README complets"),
    ("Observabilité :","Métriques Prometheus dans chaque service, dashboards Grafana infra + sécurité, ILM Elasticsearch automatique"),
]
by16=y16+Inches(0.05)
for bt,rt in bnf:
    bullet(s16,ML,by16,SW-2*ML,bt,rt,12.5); by16+=Inches(0.62)
footer(s16,16)

# ══════════════════════════════════════════════════════
# S17 — Résultats obtenus
# ══════════════════════════════════════════════════════
s17=prs.slides.add_slide(BLK)
nav_bar(s17,3)
y17=sec_title(s17,"Résultats obtenus")
res=[("🛡","9 TYPES D'ATTAQUES\nDÉTECTÉES EN TEMPS RÉEL"),
     ("📊","DASHBOARD COMPLET\nSPA + GRAFANA + KIBANA + GÉO MAP"),
     ("🔔","11 RÈGLES D'ALERTE\nPROMETHEUS + REDIS PERSISTANT")]
rw17=(SW-2*ML-Inches(0.4))/3; rx17=ML; ry17=y17+Inches(0.5)
for ico17,lbl17 in res:
    cz=Inches(1.3); cx17=rx17+rw17/2-cz/2
    sc=s17.shapes.add_shape(9,cx17,ry17,cz,cz)
    sc.fill.solid(); sc.fill.fore_color.rgb=N; sc.line.fill.background()
    txt(s17,cx17,ry17,cz,cz,ico17,38,False,W,PP_ALIGN.CENTER)
    txt(s17,rx17,ry17+cz+Inches(0.25),rw17,Inches(1.0),lbl17,12,True,DK,PP_ALIGN.CENTER,True)
    rx17+=rw17+Inches(0.2)
footer(s17,17)

# ══════════════════════════════════════════════════════
# S18 — Démo
# ══════════════════════════════════════════════════════
s18=prs.slides.add_slide(BLK)
nav_bar(s18,2)
y18=sec_title(s18,"Démonstration")
demos=[("dashboard_review_1_main.png","Dashboard principal — vue générale"),
       ("dashboard_review_3_alerts.png","Alertes de sécurité en temps réel"),
       ("dashboard_review_4_geomap.png","Carte géographique des attaques")]
dw=(SW-2*ML-Inches(0.4))/3; dh=SH-FH-y18-Inches(0.5); dx=ML
for fn,dl in demos:
    rct(s18,dx,y18+Inches(0.05),dw,dh,fc=LB,lc=GR)
    txt(s18,dx,y18+Inches(0.05),dw,dh-Inches(0.35),"📸\n["+fn+"]",11,False,GR,PP_ALIGN.CENTER)
    txt(s18,dx,y18+dh-Inches(0.28),dw,Inches(0.28),dl,9,False,RGBColor(0x77,0x77,0x77),PP_ALIGN.CENTER)
    dx+=dw+Inches(0.2)
footer(s18,18)

# ══════════════════════════════════════════════════════
# S19 — Conclusion
# ══════════════════════════════════════════════════════
s19=prs.slides.add_slide(BLK)
nav_bar(s19,4)
y19=sec_title(s19,"Conclusion")
cards_4(s19,[
    ("🎯","PFA enrichissant :","Conception et déploiement d'une plateforme SOC complète de A à Z"),
    ("💻","Compétences renforcées :","ELK Stack, Docker, Python/Flask, Prometheus/Grafana, sécurité réseau & JWT"),
    ("🤝","Architecture solide :","14 conteneurs microservices, TLS/HTTPS, JWT HS256, RBAC 3 rôles, Redis persistant"),
    ("🚀","Livrable complet :","Dashboard fonctionnel, 11 alertes automatiques, 9 scénarios d'attaque simulés"),
],y19+Inches(0.3))
footer(s19,19)

# ══════════════════════════════════════════════════════
# S20 — Perspectives
# ══════════════════════════════════════════════════════
s20=prs.slides.add_slide(BLK)
nav_bar(s20,4)
y20=sec_title(s20,"Perspectives")
cards_4(s20,[
    ("🧠","Machine Learning :","Détection d'anomalies par IA pour un scoring automatique des menaces en temps réel"),
    ("☁️","Cloud Public :","Déploiement sur AWS/Azure avec haute disponibilité, auto-scaling et Kubernetes"),
    ("🔗","SIEM Externe :","Connexion à Splunk ou QRadar pour corrélation d'événements multi-sources"),
    ("📱","Application Mobile :","App de monitoring avec notifications push pour les alertes CRITICAL"),
],y20+Inches(0.3))
footer(s20,20)

# ══════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════
out=r"C:\Users\Hicham\Desktop\pfa 2025-2026\presentation\SecureWatch_PFA_2025-2026.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
